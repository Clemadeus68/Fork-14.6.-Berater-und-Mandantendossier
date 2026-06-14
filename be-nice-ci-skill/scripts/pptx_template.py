"""
be nice network — PowerPoint-Vorlage (.pptx)
CI-konforme Präsentationsvorlage für Clemens Gutmann
Google-Slides-kompatibel.

VERWENDUNG:
  1. METADATEN unten anpassen (Titel, Untertitel, Datum)
  2. FOLIEN unten anpassen (slides_data Liste)
  3. FOOTER-LOGOS aktivieren/deaktivieren
  4. Ausführen: python3 pptx_template.py

VORAUSSETZUNGEN:
  pip install python-pptx Pillow --break-system-packages
"""

from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from PIL import Image
import io, os

# ── Logo-Pfade ─────────────────────────────────────────────────────────────────
LOGOS = os.path.join(os.path.dirname(__file__), '..', 'assets')

LOGO_BENICE  = os.path.join(LOGOS, 'logo.png')  # immer aktiv

# FOOTER-LOGOS: gewünschte aktivieren (Kommentar entfernen)
# LOGO_ACN    = os.path.join(LOGOS, 'alchimedus-netzwerk.png')
# LOGO_MASTER = os.path.join(LOGOS, 'alchimedus-master.jpg')
# LOGO_INQA   = os.path.join(LOGOS, 'inqa.jpg')

# Liste der aktiven Footer-Logos (Pfad, Ratio width/height)
# Reihenfolge = Reihenfolge auf der Folie (links → rechts in Dritteln)
FOOTER_LOGOS = [
    # (LOGO_ACN,    1550/360),
    # (LOGO_MASTER, 3661/1949),
    # (LOGO_INQA,   1.0),
]

# ── Farben ─────────────────────────────────────────────────────────────────────
GRUEN        = RGBColor(0x8C, 0xC6, 0x3E)
DUNKELGRAU   = RGBColor(0x45, 0x45, 0x44)
TUERKIS      = RGBColor(0x33, 0xAB, 0x97)
BLAU_TUERKIS = RGBColor(0x13, 0xA1, 0xB6)
HELLGRAU     = RGBColor(0xE9, 0xE9, 0xE9)
WEISS        = RGBColor(0xFF, 0xFF, 0xFF)
FAST_SCHWARZ = RGBColor(0x11, 0x11, 0x11)
MITTELGRAU   = RGBColor(0x77, 0x77, 0x77)

# ── Foliendimensionen 16:9 ─────────────────────────────────────────────────────
SW = Cm(33.867)
SH = Cm(19.05)

HEADER_H  = Cm(1.6)
FOOTER_H  = Cm(2.2)
FOOTER_Y  = SH - FOOTER_H
CONTENT_TOP = HEADER_H + Cm(0.25)
MARGIN_L  = Cm(1.5)
CONTENT_W = SW - Cm(3.0)

# ══════════════════════════════════════════════════════════════════════════════
# METADATEN — hier anpassen
# ══════════════════════════════════════════════════════════════════════════════
META = {
    "ausgabe":    "Praesentation_Titel.pptx",
    "titel":      "Präsentationstitel",
    "untertitel": "Untertitel oder Anlass",
    "datum":      "Clemens Gutmann · be nice network · Mai 2026",
}

# ── Hilfsfunktionen ────────────────────────────────────────────────────────────

def png_buf(path):
    """Konvertiert beliebiges Bild zu PNG-BytesIO (für pptx)."""
    img = Image.open(path)
    if img.mode == 'P':    img = img.convert('RGBA')
    elif img.mode == 'RGB': img = img.convert('RGBA')
    buf = io.BytesIO()
    img.save(buf, format='PNG')
    buf.seek(0)
    return buf

def img_ratio(path):
    img = Image.open(path)
    return img.size[0] / img.size[1]

def rect(slide, x, y, w, h, fill=None):
    s = slide.shapes.add_shape(1, int(x), int(y), int(w), int(h))
    s.line.fill.background()
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s

def textbox(slide, text, x, y, w, h, size,
            bold=False, italic=False, color=FAST_SCHWARZ,
            align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color
    return tb

def picture(slide, path, x, y, w, h):
    slide.shapes.add_picture(png_buf(path), int(x), int(y), int(w), int(h))

def slide_number(slide, cur, total):
    textbox(slide, f"{cur} / {total}",
            SW - Cm(2.8), SH - Cm(0.65), Cm(2.6), Cm(0.55),
            size=8, color=MITTELGRAU, align=PP_ALIGN.RIGHT)

def add_header(slide):
    """HELLGRAU-Balken, be nice Logo links, URL rechts, grüne Linie."""
    rect(slide, 0, 0, SW, HEADER_H, fill=HELLGRAU)
    rect(slide, 0, HEADER_H, SW, Pt(3), fill=GRUEN)
    lh = HEADER_H - Cm(0.2)
    lw = lh * img_ratio(LOGO_BENICE)
    picture(slide, LOGO_BENICE, Cm(0.4), Cm(0.1), lw, lh)
    url_h = Cm(0.6)
    url_y = (HEADER_H - url_h) / 2
    textbox(slide, "www.nice-network.de",
            SW - Cm(5.5), url_y, Cm(5.2), url_h,
            size=10, color=MITTELGRAU, align=PP_ALIGN.RIGHT)

def add_footer(slide):
    """Footer mit aktiven Logos in Dritteln, Trennlinie."""
    if not FOOTER_LOGOS and True:  # immer Footer-Fläche zeichnen
        rect(slide, 0, FOOTER_Y, SW, FOOTER_H, fill=RGBColor(0xF7,0xF7,0xF7))
        rect(slide, 0, FOOTER_Y, SW, Pt(1.5), fill=HELLGRAU)

    n = len(FOOTER_LOGOS)
    if n == 0:
        return

    logo_h = Cm(1.62)
    logo_y = FOOTER_Y + (FOOTER_H - logo_h) / 2

    if n <= 3:
        # Logos gleichmäßig in Dritteln verteilen
        third = SW / 3
        for i, (path, ratio) in enumerate(FOOTER_LOGOS):
            max_w = third - Cm(1.0)
            calc_w = logo_h * ratio
            lw = min(calc_w, max_w)
            lh = lw / ratio if calc_w > max_w else logo_h
            cx = third * i + third / 2
            cy = logo_y + logo_h / 2
            picture(slide, path, cx - lw/2, cy - lh/2, lw, lh)
    else:
        # Mehr als 3: gleichmäßig über volle Breite
        slot = SW / n
        for i, (path, ratio) in enumerate(FOOTER_LOGOS):
            max_w = slot - Cm(0.8)
            calc_w = logo_h * ratio
            lw = min(calc_w, max_w)
            lh = lw / ratio if calc_w > max_w else logo_h
            cx = slot * i + slot / 2
            cy = logo_y + logo_h / 2
            picture(slide, path, cx - lw/2, cy - lh/2, lw, lh)

def add_left_stripe(slide):
    """Grüner Akzentstreifen links im Content-Bereich."""
    rect(slide, 0, HEADER_H + Pt(3), Cm(0.5), FOOTER_Y - HEADER_H - Pt(3), fill=GRUEN)

# ══════════════════════════════════════════════════════════════════════════════
# FOLIEN — hier anpassen
#
# Folien-Typen:
#   "title"           — Titelfolie (weißer Hintergrund, großer Titel)
#   "section_divider" — Kapiteltrennfolie (dunkler Hintergrund)
#   "content"         — Inhaltsfolie mit Titel + Bullet-Liste
#   "quote"           — Zitat/Aussage auf türkisem Hintergrund
#
# Format je Folie:
#   ("typ", "Titel oder Kap-Nr", "Untertitel oder Haupttext", ["Bullet 1", "Bullet 2", ...])
#
# Bei "title":           field[1]=Titel, field[2]=Untertitel, field[3]=[Metazeile]
# Bei "section_divider": field[1]=Nummer ("01"), field[2]=Kapitelüberschrift, field[3]=[]
# Bei "content":         field[1]=Folientitel, field[2]="" (leer), field[3]=[Bullets]
# Bei "quote":           field[1]="" (leer), field[2]=Zitatstext, field[3]=[]
# ══════════════════════════════════════════════════════════════════════════════

slides_data = [

    ("title",
     META["titel"],
     META["untertitel"],
     [META["datum"]]),

    ("section_divider", "01", "Erster Abschnitt", []),

    ("content", "Folientitel", "", [
        "Erster Bullet-Punkt.",
        "Zweiter Bullet-Punkt.",
        "Dritter Bullet-Punkt.",
        "Vierter Bullet-Punkt.",
    ]),

    ("quote", "", "Kernaussage oder Leitgedanke, der für sich stehen soll.", []),

    ("section_divider", "02", "Zweiter Abschnitt", []),

    ("content", "Weiterer Folientitel", "", [
        "Erster Punkt.",
        "Zweiter Punkt.",
        "Dritter Punkt.",
    ]),

    ("content", "Nächster Schritt", "", [
        "Konkreter nächster Schritt.",
        "Clemens Gutmann · clemens@nice-network.de",
    ]),

]

# ── Präsentation bauen ─────────────────────────────────────────────────────────
TOTAL = len(slides_data)
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]

for idx, (stype, title, body, bullets) in enumerate(slides_data, 1):
    slide = prs.slides.add_slide(blank)

    if stype == "title":
        rect(slide, 0, 0, SW, SH, fill=WEISS)
        add_header(slide)
        add_left_stripe(slide)
        textbox(slide, title,
                Cm(1.5), CONTENT_TOP + Cm(1.8), CONTENT_W, Cm(2.8),
                size=36, bold=True, color=DUNKELGRAU)
        textbox(slide, body,
                Cm(1.5), CONTENT_TOP + Cm(4.8), CONTENT_W, Cm(1.5),
                size=22, color=TUERKIS)
        if bullets:
            textbox(slide, bullets[0],
                    Cm(1.5), CONTENT_TOP + Cm(6.3), CONTENT_W, Cm(0.9),
                    size=14, color=MITTELGRAU)
        add_footer(slide)

    elif stype == "section_divider":
        rect(slide, 0, 0, SW, SH, fill=DUNKELGRAU)
        lh = HEADER_H - Cm(0.2)
        lw = lh * img_ratio(LOGO_BENICE)
        picture(slide, LOGO_BENICE, Cm(0.4), Cm(0.1), lw, lh)
        rect(slide, 0, HEADER_H, SW, Pt(3), fill=GRUEN)
        textbox(slide, title,
                MARGIN_L, SH/2 - Cm(3.2), Cm(3), Cm(1.1),
                size=14, bold=True, color=GRUEN)
        textbox(slide, body,
                MARGIN_L, SH/2 - Cm(2.0), CONTENT_W, Cm(2.8),
                size=38, bold=True, color=WEISS)
        rect(slide, MARGIN_L, SH/2 + Cm(1.0), Cm(12), Pt(3), fill=GRUEN)
        add_footer(slide)

    elif stype == "quote":
        rect(slide, 0, 0, SW, SH, fill=TUERKIS)
        lh = HEADER_H - Cm(0.2)
        lw = lh * img_ratio(LOGO_BENICE)
        picture(slide, LOGO_BENICE, Cm(0.4), Cm(0.1), lw, lh)
        rect(slide, 0, HEADER_H, SW, Pt(3), fill=WEISS)
        textbox(slide, body,
                Cm(3.0), SH/2 - Cm(2.2), SW - Cm(6.0), Cm(4.4),
                size=22, italic=True, bold=True, color=WEISS,
                align=PP_ALIGN.CENTER, wrap=True)
        add_footer(slide)

    elif stype == "content":
        rect(slide, 0, 0, SW, SH, fill=WEISS)
        add_header(slide)
        add_left_stripe(slide)
        textbox(slide, title,
                Cm(1.5), CONTENT_TOP + Cm(0.2), CONTENT_W, Cm(1.15),
                size=24, bold=True, color=DUNKELGRAU)
        rect(slide, Cm(1.5), CONTENT_TOP + Cm(1.45), CONTENT_W, Pt(2), fill=GRUEN)
        by = CONTENT_TOP + Cm(1.75)
        for bullet in bullets:
            rect(slide, Cm(1.65), by + Cm(0.38), Cm(0.18), Cm(0.18), fill=GRUEN)
            textbox(slide, bullet,
                    Cm(2.1), by, CONTENT_W - Cm(0.7), Cm(1.4),
                    size=16, color=FAST_SCHWARZ, wrap=True)
            by += Cm(1.6)
        add_footer(slide)

    slide_number(slide, idx, TOTAL)

prs.save(META["ausgabe"])
print(f"Gespeichert: {META['ausgabe']}")
